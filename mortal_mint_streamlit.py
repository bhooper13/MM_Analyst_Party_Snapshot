import streamlit as st
import pandas as pd
import networkx as nx
import matplotlib.pyplot as plt
import mpld3
import streamlit.components.v1 as components

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT as PP_ALIGN

G = nx.read_gml('synthetic_mm.gml')
val_map = {'cryptocurrency_address': 0,
'cryptocurrency_cluster':.10,
'location':.20,
'ip_address':.3,
'osints':.4,
'party':.5,
'address':.6,
'electronic_address':.7,
'identification':.8,
'name':.9,
'phone_number':1
# 'account':,
# 'activity':,
# 'cargo':,
# 'vessel':,
}

#todo make a function that takes in a name and returns a node id and replace character name references with it

# Define function to get the number of connected nodes and names of characters most connected to the input character
def get_connection_info(name):
    #todo handle when name is not in graph
    try:
        node_id= [x for x, y in G.nodes(data=True) if y['name'] == name][0]
        if node_id in G:
            # Get the number of connected nodes
            num_connections = len(G[node_id])
            # Get the names of characters most connected to the input character
            connected_characters = sorted(G[node_id], key=lambda x: G.degree(x), reverse=True)[:5]
            connected_names = [G.nodes[c]['name'] for c in connected_characters if G.nodes[c]['name'] != 'nan']

            return num_connections, connected_names
        else:
            return None, None
    except:
        return None, None

def get_connected_nodes(name, graph):
    connected_nodes = {}
    node_id = [x for x, y in G.nodes(data=True) if y['name'] == name][0]
    # Get the neighbors of the selected character
    neighbors = list(graph.neighbors(node_id))
    # Iterate through the neighbors to find their neighbors
    for neighbor in neighbors:
        # Get the neighbors of the neighbor
        sub_neighbors = list(graph.neighbors(neighbor))
        sub_neighbor_names = [graph.nodes[node]['name'] for node in sub_neighbors if graph.nodes[node]['name'] != 'nan']
        connected_nodes[graph.nodes[neighbor]['name']] = sub_neighbor_names
    return connected_nodes

# input a value and attribute to search for matching nodes where the attribute equals the value
#todo handle multiple return node id values
#todo handle no return node id values
def search_attribute(search_value, attribute, g):
    node_id = [x for x, y in g.nodes(data=True) if y[attribute] == search_value][0]
    return node_id

def get_related_phone_numbers(node_id, g):
    try:
        related_nodes = g.neighbors(node_id)
        phone_numbers = [g.nodes[node]["phone_number"] for node in related_nodes if g.nodes[node]['phone_number'] != 'nan']
        return phone_numbers
    except:
        return None

def get_related_emails(node_id, g):
    try:
        related_nodes = g.neighbors(node_id)
        emails = [g.nodes[node]["electronic_address"] for node in related_nodes if g.nodes[node]["electronic_address"] != 'nan']
        return emails
    except:
        return None

def get_connected_nodes_with_cryptoaddress(node_id, graph):
    connected_nodes_with_address = {}
    # Get the neighbors of the selected character
    # node_id = search_attribute(name, 'name', graph)
    neighbors = list(graph.neighbors(node_id))
    # Iterate through the neighbors to find their neighbors

    crypto_addresses=[graph.nodes[neighbor]['address'] for neighbor in neighbors
                                      if graph.nodes[neighbor]['address']!= 'nan']
        # Add the sub-neighbors with 'address' attributes to the dictionary
    return crypto_addresses

def get_related_street_address(node_id, g):
    try:
        related_nodes = g.neighbors(node_id)
        emails = [g.nodes[node]["full_address"] for node in related_nodes if g.nodes[node]["full_address"] != 'nan']
        return emails
    except:
        return None

    #todo: change slide to be wider
    #todo: change the font so it's smaller

#find neighbors K number of edges away
def k_neighbors(G, node, k):
    subgraph = nx.ego_graph(G, node, radius=k)
    neighbors = list(subgraph.nodes())
    return neighbors

#find connected tables given a node id, a table name, a specified number of edges away, and a graph
def find_connected_tables(node_id, table_name, num_edges_away, graph):
    connected_tables = {}
    # Get the neighbors of the selected character
    neighbors = list(graph.neighbors(node_id))
    # Iterate through the neighbors to find their neighbors
    for neighbor in neighbors:
        # Get the neighbors of the neighbor
        sub_neighbors = list(graph.neighbors(neighbor))
        # Remove the selected character from the list of sub-neighbors
        # Get the 'name' attribute of each sub-neighbor
        sub_neighbor_names = [graph.nodes[node]['name'] for node in sub_neighbors if graph.nodes[node]['name'] != 'nan']
        # Add the sub-neighbors to the dictionary
        connected_tables[graph.nodes[neighbor]['name']] = sub_neighbor_names
    return connected_tables

def export_to_powerpoint(search_name, num_connections,related_phone_numbers, related_crypto_addresses,
                         connected_character_names, related_emails, related_street_address):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = f"{search_name} Connection Information"
    subtitle = slide.placeholders[1]
    subtitle.text = f"{search_name} is connected to {num_connections} other nodes."

    #todo: make separate bullet point sections for each result
    bullet_points = slide.shapes.placeholders[1]
    #todo: add titles
    #todo: make the bullet points smaller
    if connected_character_names:
        section_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(2.0 + 1*0.75), width=Inches(9), height=Inches(0.25))
        section_header = section_box.text_frame.add_paragraph()
        section_header.text = "Connected Entities"
        section_header.font.bold = True
        bullet_points = section_box.text_frame.add_paragraph()
        for i, name in enumerate(connected_character_names):
            p = bullet_points.add_run()
            p.text = f"{name}"+"\n"
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
    # if related_phone_numbers:
    #     for i, number in enumerate(related_phone_numbers):
    #         p = bullet_points.text_frame.add_paragraph()
    #         p.text = f"{number}"
    #         p.font.size = Pt(18)
    #         p.font.name = 'Calibri'
    if related_phone_numbers:
        section_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(2.0 + 2*0.75), width=Inches(9), height=Inches(0.25))
        section_header = section_box.text_frame.add_paragraph()
        section_header.text = "Phone Numbers"
        section_header.font.bold = True
        bullet_points = section_box.text_frame.add_paragraph()
        for i, number in enumerate(related_phone_numbers):
            p = bullet_points.add_run()
            p.text = f"{number}"+"\n"
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
    # if related_emails:
    #     for i, email in enumerate(related_emails):
    #         p = bullet_points.text_frame.add_paragraph()
    #         p.text = f"{email}"
    #         p.font.size = Pt(18)
    #         p.font.name = 'Calibri'
    if related_emails:
        section_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(2.0 + 3*0.75), width=Inches(9), height=Inches(0.25))
        section_header = section_box.text_frame.add_paragraph()
        section_header.text = "Emails"
        section_header.font.bold = True
        bullet_points = section_box.text_frame.add_paragraph()
        for i, email in enumerate(related_emails):
            p = bullet_points.add_run()
            p.text = f"{email}"+"\n"
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
    # if related_crypto_addresses:
    #     for i, crypto_address in enumerate(related_crypto_addresses):
    #         p = bullet_points.text_frame.add_paragraph()
    #         p.text = f"{crypto_address}"
    #         p.font.size = Pt(18)
    #         p.font.name = 'Calibri'
    if related_crypto_addresses:
        section_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(2.0 + 4*0.75), width=Inches(9), height=Inches(0.25))
        section_header = section_box.text_frame.add_paragraph()
        section_header.text = "Crypto Addresses"
        section_header.font.bold = True
        bullet_points = section_box.text_frame.add_paragraph()
        for i, crypto in enumerate(related_crypto_addresses):
            p = bullet_points.add_run()
            p.text = f"{crypto}"+"\n"
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
    # if related_street_address:
    #     for i, street_address in enumerate(related_street_address):
    #         p = bullet_points.text_frame.add_paragraph()
    #         p.text = f"{street_address}"
    #         p.font.size = Pt(18)
    #         p.font.name = 'Calibri'
    if related_street_address:
        section_box = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(2.0 + 5*0.75), width=Inches(9), height=Inches(0.25))
        section_header = section_box.text_frame.add_paragraph()
        section_header.text = "Street Addresses"
        section_header.font.bold = True
        bullet_points = section_box.text_frame.add_paragraph()
        for i, street_address in enumerate(related_street_address):
            p = bullet_points.add_run()
            p.text = f"{street_address}"+"\n"
            p.font.size = Pt(18)
            p.font.name = 'Calibri'
    prs.save(f"{search_name}_connections.pptx")

#streamlit functions
def list_related_information(input_list, attribute):
    if input_list:
        st.write(f"Related {attribute} are:")
        for i, value in enumerate(input_list):
            st.write(f"{i + 1}. {value}")


# Set up the Streamlit app
st.title("Mortal Mint Network Graph Query")

# Get user input for character name
search_name = st.text_input("Enter the name of a entity to search:", value ="")

#draw the full graph if no character name is entered
if search_name== '':
    # Draw the network graph of the input character and its connections
    color_map = [val_map.get(attributes['table'], 0.25) for node, attributes in G.nodes(data=True)]

    #dictionary with nodes as keys and table as the value
    node_table_labels = {node: attributes['table'] for node, attributes in G.nodes(data=True)}

    # pos = nx.spring_layout(G, k=.40, iterations=20, seed=42)
    pos = nx.kamada_kawai_layout(G)
    fig, ax = plt.subplots(figsize=(8, 8))
    nx.draw(G, pos=pos, ax=ax, with_labels=True, labels = node_table_labels, font_size=8, cmap=plt.get_cmap('viridis'),
            node_color=color_map, edge_color='gray')
    ax.set_title(f"Full Network Graph")
    fig_html = mpld3.fig_to_html(fig, figid = 'fig1')
    components.html(fig_html, height=1000, width=1000)
    # st.pyplot(fig)
    st.stop()

if search_name != "":

    # Call the function to get connection information for the input character
    num_connections, connected_character_names = get_connection_info(search_name)
    # parties_2_edges_away = get_connected_nodes(character_name, G)
    result_node_id = search_attribute(search_name, 'name', G)
    # todo: get node id and pass to these functions
    related_crypto_addresses = get_connected_nodes_with_cryptoaddress(result_node_id, G)
    related_phone_numbers=get_related_phone_numbers(result_node_id, G)
    related_emails=get_related_emails(result_node_id, G)
    related_street_address= get_related_street_address(result_node_id, G)
    # Display the connection information

if num_connections is not None and connected_character_names is not None:
    st.write(f"{search_name} is connected to {num_connections} other nodes.")
    st.write(f"The entities most connected to {search_name} are:")
    for i, name in enumerate(connected_character_names):
        st.write(f"{i + 1}. {name}")

    # node_id = [x for x, y in G.nodes(data=True) if y['name'] == name][0]

    # Draw the network graph of the input character and its connections
    subgraph = nx.ego_graph(G, result_node_id, radius=1)

    color_map = [val_map.get(attributes['table'], 0.25) for node, attributes in subgraph.nodes(data=True)]

    #dictionary with nodes as keys and table as the value
    # node_table_labels = {node: attributes['table'] + '\n' + attributes['guid']
    #                      for node, attributes in subgraph.nodes(data=True)}
    node_table_labels = {node: attributes['table']
                         for node, attributes in subgraph.nodes(data=True)}

    # pos = nx.spring_layout(subgraph, k=0.15, iterations=20, seed=42)
    pos = nx.shell_layout(subgraph, center = (0,0), nlist = [[result_node_id], [x for x in subgraph.nodes() if x != result_node_id]])
    plt.axis('off')
    fig, ax = plt.subplots(figsize=(8, 8))

    nx.draw(subgraph, pos=pos, ax=ax, with_labels=True, labels = node_table_labels,
            font_size=8, cmap=plt.get_cmap('viridis'),
            node_color=color_map, edge_color='gray')
    ax.set_title(f"Network Graph for {search_name} and Their Connections")
    # st.pyplot(fig)
    fig_html = mpld3.fig_to_html(fig)
    components.html(fig_html, height=800, width=800)
    st.title("Related Information")
    list_related_information(related_crypto_addresses, 'cryptocurrency addresses')
    list_related_information(related_phone_numbers, 'phone numbers')
    list_related_information(related_emails, 'emails')
    list_related_information(related_street_address, 'street addresses')

    # Add a button to export results to a PowerPoint slide
    if st.button("Export Results to PowerPoint"):
        export_to_powerpoint(search_name, num_connections, related_phone_numbers,
                             related_crypto_addresses,connected_character_names, related_emails, related_street_address)
    # if parties_2_edges_away:
    #     st.write(f"The characters 2 edges away from {character_name} are:")
    #     for i, party in enumerate(parties_2_edges_away.items()):
    #         st.write(f"{i + 1}. {party[0]}")
else:
    st.write("Please enter a valid entity name.")
