import streamlit as st
import networkx as nx
import matplotlib.pyplot as plt
import tempfile
import pptx
from pptx.util import Inches

def save_graph_to_ppt(graph, ppt_file):
    # Draw the graph and save it as an image
    plt.figure()
    pos = nx.spring_layout(graph)
    nx.draw(graph, pos, with_labels=True, node_size=2000)
    # nx.draw_networkx_edge_labels(graph, pos)

    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as temp_file:
        plt.savefig(temp_file.name, format='PNG', dpi=300)
        temp_file.flush()

        # Create a PowerPoint presentation or load an existing one
        if ppt_file:
            presentation = pptx.Presentation(ppt_file)
        else:
            presentation = pptx.Presentation()

        # Add a new slide to the presentation
        slide_layout = presentation.slide_layouts[5]
        slide = presentation.slides.add_slide(slide_layout)

        # Add the graph image to the slide
        slide.shapes.add_picture(temp_file.name, Inches(0), Inches(1), width=Inches(10))

        # Save the presentation
        presentation.save(ppt_file if ppt_file else 'graph_presentation.pptx')


