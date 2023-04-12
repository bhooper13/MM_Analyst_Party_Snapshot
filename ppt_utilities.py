# import streamlit as st
# import networkx as nx
# import matplotlib.pyplot as plt
# import tempfile
# import plotly.express as px
import plotly.io as pio
from pptx import Presentation
from io import BytesIO
from pptx.util import Inches, Pt
import network_utilities as nx_util
from PIL import Image
pio.kaleido.scope.mathjax = None
import os
# import psutil
# import subprocess


# def close_specific_ppt_file(file_path):
#     """
#     Closes a specific PowerPoint file if it is already open.
#
#     :param file_path: string, path to the PowerPoint file
#     """
#     # Get the file name from the file path
#     file_name = os.path.basename(file_path)
#
#     # Check if PowerPoint is running
#     ppt_process = [process for process in psutil.process_iter() if process.name() == "POWERPNT.EXE"]
#
#     if ppt_process:
#         # Get the window title of the open PowerPoint presentations
#         window_titles = subprocess.check_output(["tasklist", "/v", "/fi", "imagename eq POWERPNT.EXE"]).decode("utf-8")
#
#         # Iterate through the window titles and close the specific file if it's open
#         for title in window_titles.split("\n"):
#             if file_name in title:
#                 # Get the process ID of the PowerPoint instance with the specific file open
#                 process_id = int(title.split()[1])
#
#                 # Close the specific PowerPoint instance
#                 os.system(f"taskkill /pid {process_id}")
#
# def open_specific_ppt_file(file_path):
#     """
#     Opens a specific PowerPoint file.
#
#     :param file_path: string, path to the PowerPoint file
#     """
#     if os.path.exists(file_path):
#         subprocess.Popen(["start", "powerpnt.exe", file_path], shell=True)

def list_to_string_with_newlines(input_list):
    """
    Given a list, this function returns a string with each element separated by a newline character.

    :param input_list: list of elements
    :return: string with each element separated by a newline
    """
    return "\n".join(str(element) for element in input_list) +"\n\n"

def create_section_text(graph, attributes_name_list):
    section_text = ""
    for attribute in attributes_name_list:
        value_list = nx_util.get_attribute_values(graph, attribute)
        if value_list:
            section_text += attribute.replace('_',' ').title() +"\n"
            section_text += list_to_string_with_newlines(value_list)
    return section_text

def save_graph_to_ppt(figure, ppt_file='graph_presentation.pptx'):
    # Convert the Plotly figure to an image (PNG format)
    image_bytes = pio.to_image(figure, format="png")

    # Create a new PowerPoint presentation
    presentation = Presentation()

    # Add a slide with a title and a content placeholder
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)

    # Add a title to the slide
    title = slide.shapes.title
    title.text = "Network Graph"

    # Add the Plotly figure (image) to the slide
    image_placeholder = slide.placeholders[1]  # Placeholder for the content
    image_stream = BytesIO(image_bytes)
    slide_image = image_placeholder.insert_picture(image_stream)

    # Save the PowerPoint presentation to a file
    presentation.save(ppt_file)




def export_to_powerpoint(graph, search_value):
    prs = Presentation('MM_template.pptx')
    slide = prs.slides.add_slide(prs.slide_layouts[0])

    # for shape in slide.placeholders:
    #     print('%d %s' % (shape.placeholder_format.idx, shape.name))
    # print('*'*50)
    title = slide.shapes.title
    title.text = f"{search_value.title()} Connection Information"
    # subtitle = slide.placeholders[1]
    # subtitle.text = f"{search_value.title()} is connected to xxx other nodes."

    #get attribute values
    section_attributes = ['full_address', 'phone_number','electronic_address','party_full_name', 'address',
                          'citizenship', 'date_of_birth','party_identification_number', 'account_number',
                          'account_maximum_value_amount','sanction_type', 'sdafsdf']
    section_text = create_section_text(graph, section_attributes)

    #add text to slide
    text_box = slide.placeholders[10]
    text_box.text = (section_text).title()


    #exporting the graph as an image on a new slide
    image_slide = prs.slides.add_slide(prs.slide_layouts[1])
    image_slide_title = image_slide.shapes.title
    image_slide_title.text = f"{search_value.title()} Network Graph"
    shapes = image_slide.shapes

    # insert logo image
    left = Inches(3.725)
    top = Inches(1)
    image_picture = shapes.add_picture('figure.png', left, top)

    # subtitle = image_slide.shapes[1]
    # sp = subtitle.element
    # sp.getparent().remove(sp)
    prs.save(f"{search_value}_mortal_mint_snapshot.pptx")
    # open_specific_ppt_file(f"{search_value}_mortal_mint_snapshot.pptx")